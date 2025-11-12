from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import re
import shutil

MEDIA_EXPORT_DIR = "docs/media"

def ensure_dir(p: Path):
    p.mkdir(parents=True, exist_ok=True)

def sanitize_stem(stem: str) -> str:
    s = re.sub(r"\s+", "_", stem)
    s = re.sub(r"[^A-Za-z0-9_\-]", "", s)
    return s or "slides"

def iter_shape_text(shape):
    if hasattr(shape, "has_text_frame") and shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            level = getattr(para, "level", 0)
            text = "".join(run.text for run in para.runs) or para.text or ""
            if text and text.strip():
                yield level, text.strip()

def collect_media(prs: Presentation, pptx_path: Path, out_dir: Path, slide_index: int):
    """
    Ekstrak file media mentah dari paket PPTX (gambar) lalu buat mapping.
    Catatan: python-pptx tidak expose langsung objek image per shape untuk export individual alt text,
    jadi pendekatan sederhana: copy semua media ke folder docs/media, lalu di slide catat urutan.
    """
    # File media disimpan di ppt/media/ dalam struktur zip PPTX. Kita bisa akses manual:
    # PPTX sebenarnya ZIP → ekstraksi semua lalu ambil folder ppt/media
    temp_unzip = out_dir / "__temp_unzip__"
    if temp_unzip.exists():
        shutil.rmtree(temp_unzip)
    ensure_dir(temp_unzip)
    # Unzip
    shutil.unpack_archive(str(pptx_path), str(temp_unzip), "zip")
    media_root = temp_unzip / "ppt" / "media"
    exported = []
    if media_root.exists():
        export_target = out_dir / MEDIA_EXPORT_DIR
        ensure_dir(export_target)
        for m in media_root.iterdir():
            # Copy jika belum ada
            target = export_target / m.name
            if not target.exists():
                shutil.copy2(m, target)
            exported.append(m.name)
    # Bersihkan temp
    shutil.rmtree(temp_unzip, ignore_errors=True)
    return exported

def slide_notes(slide):
    try:
        notes = slide.notes_slide.notes_text
        return notes.strip() if notes and notes.strip() else ""
    except Exception:
        return ""

def shape_summary(slide):
    """
    Ringkas non-text shapes: gambar, tabel, chart, SmartArt dll.
    """
    counts = {
        "picture": 0,
        "table": 0,
        "chart": 0,
        "group": 0,
        "other": 0
    }
    for shape in slide.shapes:
        st = shape.shape_type
        if st == MSO_SHAPE_TYPE.PICTURE:
            counts["picture"] += 1
        elif st == MSO_SHAPE_TYPE.TABLE:
            counts["table"] += 1
        elif st == MSO_SHAPE_TYPE.CHART:
            counts["chart"] += 1
        elif st == MSO_SHAPE_TYPE.GROUP:
            counts["group"] += 1
        else:
            # Filter text-frame shapes keluar agar 'other' tidak overcount
            if not (hasattr(shape, "has_text_frame") and shape.has_text_frame):
                counts["other"] += 1
    parts = []
    for k, v in counts.items():
        if v:
            parts.append(f"{k}:{v}")
    return ", ".join(parts) if parts else "tidak ada shape non-teks"

def extract_one(pptx_path: Path, out_dir: Path):
    prs = Presentation(str(pptx_path))
    stem = sanitize_stem(pptx_path.stem)
    out_path = out_dir / f"{stem}.md"

    media_list = collect_media(prs, pptx_path, out_dir, 0)  # sekali saja

    lines = [f"# {pptx_path.name}", ""]
    lines.append(f"**Total media diekstrak (semua slide)**: {len(media_list)}")
    lines.append("")
    for i, slide in enumerate(prs.slides, start=1):
        lines.append(f"## Slide {i}")
        # Teks
        any_text = False
        for level, text in iter_shape_text_loop(slide):
            any_text = True
            indent = "  " * min(level, 4)
            lines.append(f"{indent}- {text}")
        if not any_text:
            lines.append("_(tidak ada teks terdeteksi di slide ini)_")

        # Notes
        notes = slide_notes(slide)
        if notes:
            lines.append("")
            lines.append("**Catatan (Notes):**")
            lines.append(f"> {notes}")

        # Shape summary
        summary = shape_summary(slide)
        lines.append("")
        lines.append(f"**Ringkasan shape non-teks:** {summary}")

        # Placeholder rekomendasi alt text (jika belum ada teks dan ada gambar)
        if "picture:" in summary and not any_text and not notes:
            lines.append("")
            lines.append("_Rekomendasi: Tambahkan alt text / catatan untuk gambar di slide ini._")

        lines.append("")

    out_path.write_text("\n".join(lines), encoding="utf-8")
    print(f"Wrote: {out_path}")

def iter_shape_text_loop(slide):
    for shape in slide.shapes:
        yield from iter_shape_text(shape)

def main():
    repo_root = Path(__file__).resolve().parents[1]
    out_dir = repo_root / "docs"
    ensure_dir(out_dir)

    pptx_files = sorted(repo_root.rglob("*.pptx"))
    if not pptx_files:
        print("No PPTX files found.")
        return

    for f in pptx_files:
        if "docs" in f.parts:
            continue
        extract_one(f, out_dir)

if __name__ == "__main__":
    main()
